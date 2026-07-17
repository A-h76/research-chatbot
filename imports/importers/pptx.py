class PptxImporter:
    extensions = (".pptx", ".pptm")
    supports_locators = False

    def matches(self, lower_name, mime):
        return lower_name.endswith((".pptx", ".pptm"))

    def extract(self, path, mime, name):
        from pptx import Presentation
        prs = Presentation(path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"--- Slide {i} ---")
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    parts.append(shape.text_frame.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        parts.append(" | ".join(c.text for c in row.cells))
            notes = getattr(slide, "notes_slide", None)
            if slide.has_notes_slide and notes and notes.notes_text_frame.text.strip():
                parts.append(f"[Notes] {notes.notes_text_frame.text}")
        return "\n".join(parts).strip()
